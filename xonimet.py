#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
XONIMET 2026 - Extractor Universal de Metadatos (Con Reporte PDF)
Extrae metadatos de archivos y genera reportes en PDF con formato profesional.
Optimizado para AUR - Configuracion en ~/.xonimet/
Desarrollador: Darian Alberto Camacho Salas
Organizacion: XONIDU
"""

import os
import sys
import json
import datetime
import hashlib
from pathlib import Path

# ============================================================================
# Verificacion de dependencias con mensajes claros
# ============================================================================
try:
    from PIL import Image
    from PIL.ExifTags import TAGS
except ImportError:
    print("[ERROR] Pillow no instalado. Ejecuta: pip install pillow")
    sys.exit(1)

try:
    import mutagen
    from mutagen.mp3 import MP3
    from mutagen.flac import FLAC
    from mutagen.mp4 import MP4
    from mutagen.oggvorbis import OggVorbis
except ImportError:
    print("[ERROR] Mutagen no instalado. Ejecuta: pip install mutagen")
    sys.exit(1)

try:
    import ffmpeg
except ImportError:
    print("[ERROR] ffmpeg-python no instalado. Ejecuta: pip install ffmpeg-python")
    sys.exit(1)

try:
    from PyPDF2 import PdfReader
except ImportError:
    print("[ERROR] PyPDF2 no instalado. Ejecuta: pip install PyPDF2")
    sys.exit(1)

try:
    import docx
except ImportError:
    print("[ERROR] python-docx no instalado. Ejecuta: pip install python-docx")
    sys.exit(1)

try:
    from openpyxl import load_workbook
except ImportError:
    print("[ERROR] openpyxl no instalado. Ejecuta: pip install openpyxl")
    sys.exit(1)

try:
    from pptx import Presentation
except ImportError:
    print("[ERROR] python-pptx no instalado. Ejecuta: pip install python-pptx")
    sys.exit(1)

try:
    import exifread
except ImportError:
    print("[ERROR] exifread no instalado. Ejecuta: pip install exifread")
    sys.exit(1)

try:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.lib import colors
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
    from reportlab.lib.enums import TA_CENTER
except ImportError:
    print("[ERROR] reportlab no instalado. Ejecuta: pip install reportlab")
    sys.exit(1)

import warnings
warnings.filterwarnings('ignore')

# ============================================================================
# Colores para terminal
# ============================================================================
class Colors:
    HEADER = '\033[95m'
    BLUE = '\033[94m'
    GREEN = '\033[92m'
    YELLOW = '\033[93m'
    RED = '\033[91m'
    END = '\033[0m'
    BOLD = '\033[1m'
    CYAN = '\033[96m'
    PURPLE = '\033[95m'

# ============================================================================
# Clase principal Xonimet
# ============================================================================
class Xonimet:
    def __init__(self, file_path=None):
        self.file_path = Path(file_path) if file_path else None
        self.metadata = {}
    
    def get_config_dir(self):
        """Retorna el directorio de configuracion en ~/.xonimet/"""
        return os.path.join(os.path.expanduser("~"), '.xonimet')
    
    def ensure_config_dir(self):
        """Asegura que el directorio de configuracion existe"""
        config_dir = self.get_config_dir()
        os.makedirs(config_dir, exist_ok=True)
        return config_dir
    
    def get_pdf_output_path(self, base_name=None):
        """Genera ruta de salida para PDF en ~/.xonimet/ o directorio actual"""
        if base_name is None:
            base_name = Path(self.file_path).stem if self.file_path else "reporte"
        
        config_dir = self.get_config_dir()
        # Intentar guardar en config dir primero
        if os.access(config_dir, os.W_OK):
            return os.path.join(config_dir, f"{base_name}_reporte.pdf")
        else:
            # Fallback al directorio actual
            return f"{base_name}_reporte.pdf"
    
    def set_file(self, file_path):
        """Establece el archivo a analizar"""
        self.file_path = Path(file_path)
        self.metadata = {}
    
    def extract_all(self):
        """Extrae todos los metadatos del archivo"""
        if not self.file_path or not self.file_path.exists():
            return {'error': 'Archivo no existe'}
        
        self.metadata = {
            'archivo': {
                'nombre': self.file_path.name,
                'ruta': str(self.file_path.absolute()),
                'tamaño_bytes': self.file_path.stat().st_size,
                'tamaño_formateado': self._format_bytes(self.file_path.stat().st_size),
                'creado': datetime.datetime.fromtimestamp(self.file_path.stat().st_ctime).isoformat(),
                'modificado': datetime.datetime.fromtimestamp(self.file_path.stat().st_mtime).isoformat(),
                'accedido': datetime.datetime.fromtimestamp(self.file_path.stat().st_atime).isoformat(),
                'extension': self.file_path.suffix.lower(),
                'tipo_mime': self._get_mime_type(),
                'hashes': self._calculate_hashes()
            },
            'metadatos_especificos': self._extract_specific_metadata()
        }
        return self.metadata
    
    def _format_bytes(self, bytes):
        """Formatea bytes a unidades legibles"""
        for unit in ['B', 'KB', 'MB', 'GB', 'TB']:
            if bytes < 1024.0:
                return f"{bytes:.2f} {unit}"
            bytes /= 1024.0
        return f"{bytes:.2f} PB"
    
    def _get_mime_type(self):
        """Obtiene el tipo MIME del archivo"""
        import mimetypes
        mime_type, _ = mimetypes.guess_type(str(self.file_path))
        return mime_type or 'desconocido'
    
    def _calculate_hashes(self):
        """Calcula hashes del archivo"""
        hashes = {}
        try:
            with open(self.file_path, 'rb') as f:
                data = f.read()
                hashes['md5'] = hashlib.md5(data).hexdigest()
                hashes['sha1'] = hashlib.sha1(data).hexdigest()
                hashes['sha256'] = hashlib.sha256(data).hexdigest()
        except:
            hashes['error'] = 'No se pudo calcular hashes'
        return hashes
    
    def _extract_image_metadata(self):
        """Extrae metadatos de imagenes"""
        metadata = {}
        try:
            img = Image.open(self.file_path)
            metadata.update({
                'dimensiones': f"{img.width} x {img.height}",
                'modo': img.mode,
                'formato': img.format,
                'info_basica': {
                    'ancho': img.width,
                    'alto': img.height,
                    'proporcion': round(img.width / img.height, 2)
                }
            })
            
            if hasattr(img, '_getexif') and img._getexif():
                exif = {}
                for tag_id, value in img._getexif().items():
                    tag = TAGS.get(tag_id, tag_id)
                    if isinstance(value, bytes):
                        try:
                            value = value.decode('utf-8', errors='ignore')
                        except:
                            value = str(value)
                    exif[tag] = str(value)
                metadata['exif'] = exif
            
            with open(self.file_path, 'rb') as f:
                tags = exifread.process_file(f, details=True)
                if tags:
                    metadata['exif_detallado'] = {str(k): str(v) for k, v in tags.items()}
                    
        except Exception as e:
            metadata['error'] = str(e)
        return metadata
    
    def _extract_audio_metadata(self):
        """Extrae metadatos de archivos de audio"""
        metadata = {}
        try:
            audio = mutagen.File(self.file_path)
            if audio:
                metadata['formato'] = type(audio).__name__
                metadata['duracion_segundos'] = audio.info.length
                metadata['duracion_formateado'] = str(datetime.timedelta(seconds=int(audio.info.length)))
                
                if hasattr(audio.info, 'bitrate'):
                    metadata['bitrate'] = f"{audio.info.bitrate // 1000} kbps"
                
                if hasattr(audio.info, 'sample_rate'):
                    metadata['frecuencia_muestreo'] = f"{audio.info.sample_rate} Hz"
                
                if hasattr(audio, 'tags') and audio.tags:
                    tags = {}
                    for key, value in audio.tags.items():
                        if value:
                            tags[key] = str(value[0]) if isinstance(value, list) else str(value)
                    metadata['etiquetas'] = tags
                    
        except Exception as e:
            metadata['error'] = str(e)
        return metadata
    
    def _extract_video_metadata(self):
        """Extrae metadatos de videos"""
        metadata = {}
        try:
            probe = ffmpeg.probe(str(self.file_path))
            
            if 'format' in probe:
                fmt = probe['format']
                metadata['formato'] = {
                    'nombre': fmt.get('format_name'),
                    'duracion': fmt.get('duration'),
                    'bitrate': fmt.get('bit_rate'),
                    'tamaño': fmt.get('size'),
                    'tags': fmt.get('tags', {})
                }
            
            streams = []
            for stream in probe.get('streams', []):
                stream_info = {
                    'tipo': stream.get('codec_type'),
                    'codec': stream.get('codec_name'),
                    'perfil': stream.get('profile')
                }
                
                if stream['codec_type'] == 'video':
                    stream_info.update({
                        'resolucion': f"{stream.get('width')}x{stream.get('height')}",
                        'fps': eval(stream.get('r_frame_rate', '0/1')),
                        'pixeles': stream.get('pix_fmt')
                    })
                elif stream['codec_type'] == 'audio':
                    stream_info.update({
                        'canales': stream.get('channels'),
                        'frecuencia': stream.get('sample_rate')
                    })
                
                streams.append(stream_info)
            
            metadata['streams'] = streams
            
        except Exception as e:
            metadata['error'] = str(e)
        return metadata
    
    def _extract_pdf_metadata(self):
        """Extrae metadatos de PDFs"""
        metadata = {}
        try:
            pdf = PdfReader(self.file_path)
            metadata.update({
                'paginas': len(pdf.pages),
                'encriptado': pdf.is_encrypted,
                'metadatos': dict(pdf.metadata) if pdf.metadata else {}
            })
            
            if len(pdf.pages) > 0:
                first_page = pdf.pages[0]
                text = first_page.extract_text()
                metadata['primeras_palabras'] = text[:200] + '...' if len(text) > 200 else text
                
        except Exception as e:
            metadata['error'] = str(e)
        return metadata
    
    def _extract_docx_metadata(self):
        """Extrae metadatos de documentos Word"""
        metadata = {}
        try:
            doc = docx.Document(self.file_path)
            core_props = doc.core_properties
            metadata.update({
                'autor': core_props.author,
                'creador': core_props.created,
                'modificado_por': core_props.last_modified_by,
                'fecha_modificacion': core_props.modified,
                'titulo': core_props.title,
                'asunto': core_props.subject,
                'palabras_clave': core_props.keywords,
                'categoria': core_props.category,
                'comentarios': core_props.comments,
                'parrafos': len(doc.paragraphs),
                'tablas': len(doc.tables)
            })
        except Exception as e:
            metadata['error'] = str(e)
        return metadata
    
    def _extract_xlsx_metadata(self):
        """Extrae metadatos de Excel"""
        metadata = {}
        try:
            wb = load_workbook(self.file_path, data_only=True)
            metadata.update({
                'hojas': wb.sheetnames,
                'hojas_activas': len(wb.sheetnames),
                'propiedades': {
                    'creador': wb.properties.creator,
                    'creado': str(wb.properties.created) if wb.properties.created else None,
                    'modificado': str(wb.properties.modified) if wb.properties.modified else None,
                    'titulo': wb.properties.title
                }
            })
            
            sheets_info = {}
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                sheets_info[sheet_name] = {
                    'filas': sheet.max_row,
                    'columnas': sheet.max_column,
                    'celdas_con_datos': sum(1 for row in sheet.iter_rows() for cell in row if cell.value)
                }
            metadata['detalle_hojas'] = sheets_info
            
        except Exception as e:
            metadata['error'] = str(e)
        return metadata
    
    def _extract_pptx_metadata(self):
        """Extrae metadatos de PowerPoint"""
        metadata = {}
        try:
            prs = Presentation(self.file_path)
            core_props = prs.core_properties
            metadata.update({
                'diapositivas': len(prs.slides),
                'autor': core_props.author,
                'creado': str(core_props.created) if core_props.created else None,
                'modificado': str(core_props.modified) if core_props.modified else None,
                'titulo': core_props.title,
                'asunto': core_props.subject
            })
            
            slide_stats = {'texto': 0, 'imagenes': 0, 'tablas': 0, 'graficos': 0}
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        slide_stats['texto'] += 1
                    if hasattr(shape, 'image'):
                        slide_stats['imagenes'] += 1
                    if shape.has_table:
                        slide_stats['tablas'] += 1
                    if hasattr(shape, 'chart'):
                        slide_stats['graficos'] += 1
            metadata['estadisticas_diapositivas'] = slide_stats
            
        except Exception as e:
            metadata['error'] = str(e)
        return metadata
    
    def _extract_text_metadata(self):
        """Extrae metadatos basicos de archivos de texto"""
        metadata = {}
        try:
            with open(self.file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
                lines = content.split('\n')
                metadata.update({
                    'lineas': len(lines),
                    'palabras': len(content.split()),
                    'caracteres': len(content),
                    'caracteres_sin_espacios': len(content.replace(' ', '').replace('\n', '').replace('\t', '')),
                    'primeras_10_lineas': lines[:10] if len(lines) > 10 else lines
                })
        except Exception as e:
            metadata['error'] = str(e)
        return metadata
    
    def _extract_specific_metadata(self):
        """Determina el tipo de archivo y extrae metadatos especificos"""
        ext = self.file_path.suffix.lower()
        
        if ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp', '.heic']:
            return self._extract_image_metadata()
        elif ext in ['.mp3', '.flac', '.wav', '.ogg', '.m4a', '.aac', '.wma']:
            return self._extract_audio_metadata()
        elif ext in ['.mp4', '.avi', '.mov', '.mkv', '.wmv', '.flv', '.webm', '.m4v']:
            return self._extract_video_metadata()
        elif ext == '.pdf':
            return self._extract_pdf_metadata()
        elif ext in ['.docx', '.doc']:
            return self._extract_docx_metadata()
        elif ext in ['.xlsx', '.xls']:
            return self._extract_xlsx_metadata()
        elif ext in ['.pptx', '.ppt']:
            return self._extract_pptx_metadata()
        elif ext in ['.txt', '.csv', '.json', '.xml', '.html', '.css', '.js', '.py', '.md']:
            return self._extract_text_metadata()
        else:
            return {'mensaje': 'Tipo de archivo no soportado para extraccion especifica'}
    
    def generate_pdf_report(self, output_path=None):
        """Genera un reporte PDF bonito con los metadatos"""
        if not self.metadata:
            return None
        
        # Crear directorio de configuracion si no existe
        self.ensure_config_dir()
        
        if output_path is None:
            base_name = Path(self.metadata['archivo']['nombre']).stem
            output_path = self.get_pdf_output_path(base_name)
        
        doc = SimpleDocTemplate(output_path, pagesize=A4,
                               rightMargin=72, leftMargin=72,
                               topMargin=72, bottomMargin=72)
        
        story = []
        
        # Estilos
        styles = getSampleStyleSheet()
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            textColor=colors.HexColor('#1a5490'),
            spaceAfter=30,
            alignment=TA_CENTER
        )
        heading_style = ParagraphStyle(
            'CustomHeading',
            parent=styles['Heading2'],
            fontSize=16,
            textColor=colors.HexColor('#2c7ab1'),
            spaceAfter=12,
            spaceBefore=20
        )
        subheading_style = ParagraphStyle(
            'CustomSubHeading',
            parent=styles['Heading3'],
            fontSize=14,
            textColor=colors.HexColor('#4a90c4'),
            spaceAfter=8,
            spaceBefore=12
        )
        normal_style = ParagraphStyle(
            'CustomNormal',
            parent=styles['Normal'],
            fontSize=10,
            spaceAfter=6
        )
        
        # Titulo principal
        title_text = "Reporte de Metadatos"
        story.append(Paragraph(title_text, title_style))
        story.append(Spacer(1, 0.2*inch))
        
        # Subtitulo
        subtitle_text = f"Archivo: {self.metadata['archivo']['nombre']}"
        story.append(Paragraph(subtitle_text, heading_style))
        story.append(Spacer(1, 0.2*inch))
        
        # Fecha del reporte
        fecha_text = f"Generado el: {datetime.datetime.now().strftime('%d/%m/%Y %H:%M:%S')}"
        story.append(Paragraph(fecha_text, normal_style))
        story.append(Spacer(1, 0.3*inch))
        
        # ===== INFORMACION DEL ARCHIVO =====
        story.append(Paragraph("Informacion del Archivo", heading_style))
        
        archivo_data = []
        archivo_info = self.metadata['archivo']
        for key, value in archivo_info.items():
            if key != 'hashes':
                key_name = key.replace('_', ' ').title()
                archivo_data.append([Paragraph(f"<b>{key_name}</b>", normal_style), 
                                    Paragraph(str(value), normal_style)])
        
        archivo_table = Table(archivo_data, colWidths=[2*inch, 3.5*inch])
        archivo_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (0, -1), colors.HexColor('#e8f0f8')),
            ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#c0c0c0')),
            ('VALIGN', (0, 0), (-1, -1), 'TOP'),
            ('PADDING', (0, 0), (-1, -1), 6),
        ]))
        story.append(archivo_table)
        story.append(Spacer(1, 0.2*inch))
        
        # ===== HASHES =====
        if 'hashes' in archivo_info and archivo_info['hashes']:
            story.append(Paragraph("Hashes de Seguridad", heading_style))
            
            hash_data = []
            for algo, hash_value in archivo_info['hashes'].items():
                hash_data.append([Paragraph(f"<b>{algo.upper()}</b>", normal_style),
                                Paragraph(str(hash_value), normal_style)])
            
            hash_table = Table(hash_data, colWidths=[1.2*inch, 4.3*inch])
            hash_table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (0, -1), colors.HexColor('#e8f0f8')),
                ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#c0c0c0')),
                ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                ('PADDING', (0, 0), (-1, -1), 6),
            ]))
            story.append(hash_table)
            story.append(Spacer(1, 0.2*inch))
        
        # ===== METADATOS ESPECIFICOS =====
        story.append(Paragraph("Metadatos Especificos", heading_style))
        
        spec = self.metadata.get('metadatos_especificos', {})
        
        if 'error' in spec:
            story.append(Paragraph(f"<font color='red'>Error: {spec['error']}</font>", normal_style))
        else:
            for key, value in spec.items():
                if isinstance(value, dict):
                    story.append(Paragraph(f"<b>{key.replace('_', ' ').title()}</b>", subheading_style))
                    
                    sub_data = []
                    for subkey, subvalue in value.items():
                        if subvalue:
                            subkey_name = subkey.replace('_', ' ').title()
                            sub_data.append([Paragraph(f"<i>{subkey_name}</i>", normal_style),
                                           Paragraph(str(subvalue), normal_style)])
                    
                    if sub_data:
                        sub_table = Table(sub_data, colWidths=[2*inch, 3.5*inch])
                        sub_table.setStyle(TableStyle([
                            ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#d0d0d0')),
                            ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                            ('PADDING', (0, 0), (-1, -1), 4),
                        ]))
                        story.append(sub_table)
                        story.append(Spacer(1, 0.1*inch))
                else:
                    if value:
                        line_data = [[Paragraph(f"<b>{key.replace('_', ' ').title()}</b>", normal_style),
                                     Paragraph(str(value), normal_style)]]
                        line_table = Table(line_data, colWidths=[2*inch, 3.5*inch])
                        line_table.setStyle(TableStyle([
                            ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#e0e0e0')),
                            ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                            ('PADDING', (0, 0), (-1, -1), 4),
                        ]))
                        story.append(line_table)
        
        # Pie de pagina
        story.append(Spacer(1, 0.5*inch))
        footer_text = "Reporte generado por XONIMET 2026 - Extractor Universal de Metadatos"
        story.append(Paragraph(footer_text, ParagraphStyle(
            'Footer',
            parent=styles['Normal'],
            fontSize=8,
            textColor=colors.HexColor('#888888'),
            alignment=TA_CENTER
        )))
        
        # Construir PDF
        doc.build(story)
        return output_path
    
    def print_metadata(self, metadata=None):
        """Imprime los metadatos de forma formateada"""
        if metadata is None:
            metadata = self.metadata
        
        if not metadata:
            print(f"{Colors.RED}No hay metadatos para mostrar{Colors.END}")
            return
        
        print(f"\n{Colors.CYAN}{Colors.BOLD}═══════════════════════════════════════════════════════════{Colors.END}")
        print(f"{Colors.CYAN}{Colors.BOLD}METADATOS DEL ARCHIVO{Colors.END}")
        print(f"{Colors.CYAN}{Colors.BOLD}═══════════════════════════════════════════════════════════{Colors.END}")
        
        print(f"\n{Colors.GREEN}{Colors.BOLD}INFORMACION BASICA:{Colors.END}")
        archivo = metadata.get('archivo', {})
        for key, value in archivo.items():
            if key != 'hashes':
                print(f"  {Colors.YELLOW}•{Colors.END} {key.replace('_', ' ').title()}: {value}")
        
        if 'hashes' in archivo:
            print(f"\n{Colors.GREEN}{Colors.BOLD}HASHES:{Colors.END}")
            for algo, hash_value in archivo['hashes'].items():
                print(f"  {Colors.YELLOW}•{Colors.END} {algo.upper()}: {hash_value}")
        
        if 'metadatos_especificos' in metadata and metadata['metadatos_especificos']:
            print(f"\n{Colors.GREEN}{Colors.BOLD}METADATOS ESPECIFICOS:{Colors.END}")
            spec = metadata['metadatos_especificos']
            
            if 'error' in spec:
                print(f"  {Colors.RED}⚠ {spec['error']}{Colors.END}")
            else:
                for key, value in spec.items():
                    if isinstance(value, dict):
                        print(f"\n  {Colors.CYAN}• {key.replace('_', ' ').title()}:{Colors.END}")
                        for subkey, subvalue in value.items():
                            if subvalue:
                                print(f"    {Colors.YELLOW}•{Colors.END} {subkey}: {subvalue}")
                    else:
                        print(f"  {Colors.YELLOW}•{Colors.END} {key.replace('_', ' ').title()}: {value}")
        
        print(f"\n{Colors.CYAN}{Colors.BOLD}═══════════════════════════════════════════════════════════{Colors.END}")

# ============================================================================
# Funciones del modo interactivo
# ============================================================================
def clear_screen():
    """Limpia la pantalla segun el sistema"""
    os.system('cls' if os.name == 'nt' else 'clear')

def print_menu():
    """Muestra el menu principal"""
    menu = f"""
{Colors.CYAN}{Colors.BOLD}═══════════════════════════════════════════════════════════
                   XONIMET 2026 v2.0                    
              Extractor Universal de Metadatos           
                 CON GENERADOR DE REPORTES PDF          
                    MODO INTERACTIVO                      
═══════════════════════════════════════════════════════════{Colors.END}

{Colors.GREEN}ARCHIVOS SOPORTADOS:{Colors.END}
  Imagenes | Audio | Video | PDF/DOCS | Texto

{Colors.YELLOW}═══════════════════════════════════════════════════════════{Colors.END}

{Colors.BOLD}MENU PRINCIPAL:{Colors.END}
  {Colors.CYAN}[1]{Colors.END} Seleccionar archivo para analizar
  {Colors.CYAN}[2]{Colors.END} Analizar archivo actual
  {Colors.CYAN}[3]{Colors.END} Guardar resultados en JSON
  {Colors.CYAN}[4]{Colors.END} GENERAR REPORTE PDF
  {Colors.CYAN}[5]{Colors.END} Ver informacion del archivo actual
  {Colors.CYAN}[6]{Colors.END} Cambiar archivo
  {Colors.CYAN}[7]{Colors.END} Ayuda / Formatos soportados
  {Colors.CYAN}[8]{Colors.END} Limpiar pantalla
  {Colors.CYAN}[0]{Colors.END} Salir

{Colors.YELLOW}═══════════════════════════════════════════════════════════{Colors.END}
"""
    print(menu)

def print_help():
    """Muestra ayuda detallada"""
    help_text = f"""
{Colors.CYAN}{Colors.BOLD}═══════════════════════════════════════════════════════════
AYUDA - FORMATOS SOPORTADOS
═══════════════════════════════════════════════════════════{Colors.END}

{Colors.GREEN}IMAGENES:{Colors.END}
  • .jpg, .jpeg, .png, .gif, .bmp, .tiff, .webp, .heic
  {Colors.YELLOW}→{Colors.END} EXIF, dimensiones, GPS, modelo camara, fecha

{Colors.GREEN}AUDIO:{Colors.END}
  • .mp3, .flac, .wav, .ogg, .m4a, .aac, .wma
  {Colors.YELLOW}→{Colors.END} Duracion, bitrate, etiquetas ID3, artista, album

{Colors.GREEN}VIDEO:{Colors.END}
  • .mp4, .avi, .mov, .mkv, .wmv, .flv, .webm
  {Colors.YELLOW}→{Colors.END} Resolucion, codecs, fps, streams

{Colors.GREEN}DOCUMENTOS:{Colors.END}
  • .pdf: paginas, autor, titulo, metadatos
  • .docx, .doc: autor, fechas, estadisticas
  • .xlsx, .xls: hojas, celdas, propiedades
  • .pptx, .ppt: diapositivas, estadisticas

{Colors.GREEN}TEXTO:{Colors.END}
  • .txt, .csv, .json, .xml, .html, .css, .js, .py, .md
  {Colors.YELLOW}→{Colors.END} Lineas, palabras, caracteres

{Colors.GREEN}REPORTE PDF:{Colors.END}
  • Genera un PDF con formato profesional
  • Incluye tablas, colores y estructura clara
  • Se guarda en: ~/.xonimet/ o en el directorio actual

{Colors.CYAN}{Colors.BOLD}═══════════════════════════════════════════════════════════{Colors.END}
"""
    print(help_text)

def select_file():
    """Selecciona un archivo interactivamente"""
    print(f"\n{Colors.CYAN}SELECCIONAR ARCHIVO{Colors.END}")
    print(f"{Colors.YELLOW}Escribe la ruta del archivo (o 'cancel' para volver):{Colors.END}")
    
    while True:
        file_path = input(f"{Colors.GREEN}→{Colors.END} ").strip()
        
        if file_path.lower() == 'cancel':
            return None
        
        if not file_path:
            continue
        
        file_path = os.path.expanduser(file_path)
        
        if os.path.exists(file_path):
            return file_path
        else:
            print(f"{Colors.RED}El archivo no existe. Intenta de nuevo:{Colors.END}")

def save_to_json(metadata):
    """Guarda los metadatos en un archivo JSON"""
    if not metadata:
        print(f"{Colors.RED}No hay metadatos para guardar{Colors.END}")
        return
    
    original_name = metadata.get('archivo', {}).get('nombre', 'desconocido')
    json_name = f"{Path(original_name).stem}_metadatos.json"
    
    try:
        with open(json_name, 'w', encoding='utf-8') as f:
            json.dump(metadata, f, indent=2, ensure_ascii=False, default=str)
        print(f"{Colors.GREEN}Metadatos guardados en: {json_name}{Colors.END}")
    except Exception as e:
        print(f"{Colors.RED}Error guardando: {e}{Colors.END}")

def generate_pdf(xonimet):
    """Genera reporte PDF y muestra la ruta"""
    if not xonimet.metadata:
        print(f"{Colors.RED}Primero analiza un archivo (opcion 2){Colors.END}")
        return
    
    print(f"\n{Colors.CYAN}Generando reporte PDF...{Colors.END}")
    try:
        pdf_path = xonimet.generate_pdf_report()
        # Mostrar la ruta completa del PDF generado
        print(f"{Colors.GREEN}✓ PDF generado exitosamente{Colors.END}")
        print(f"{Colors.CYAN}📄 Ruta del reporte: {os.path.abspath(pdf_path)}{Colors.END}")
        
        # Mostrar info adicional
        if os.path.exists(pdf_path):
            size = os.path.getsize(pdf_path)
            if size < 1024:
                size_str = f"{size} B"
            elif size < 1024*1024:
                size_str = f"{size/1024:.2f} KB"
            else:
                size_str = f"{size/(1024*1024):.2f} MB"
            print(f"{Colors.YELLOW}📊 Tamaño del archivo: {size_str}{Colors.END}")
    except Exception as e:
        print(f"{Colors.RED}Error generando PDF: {e}{Colors.END}")
        print(f"{Colors.YELLOW}Asegurate de tener instalado reportlab: pip install reportlab{Colors.END}")

def get_xonimet_path():
    """Detecta la ruta de xonimet.py en multiples ubicaciones (para AUR)"""
    script_dir = os.path.dirname(os.path.abspath(__file__))
    rutas = [
        os.path.join(script_dir, 'xonimet.py'),
        '/usr/share/xonimet/xonimet.py',
        os.path.join(os.path.expanduser("~"), '.xonimet', 'xonimet.py'),
        os.path.join(os.getcwd(), 'xonimet.py')
    ]
    for r in rutas:
        if os.path.exists(r):
            return r
    return None

def main():
    """Funcion principal del modo interactivo"""
    # Asegurar directorio de configuracion
    config_dir = os.path.join(os.path.expanduser("~"), '.xonimet')
    os.makedirs(config_dir, exist_ok=True)
    
    xonimet = Xonimet()
    current_file = None
    
    while True:
        clear_screen()
        print_menu()
        
        if current_file:
            print(f"{Colors.GREEN}📁 Archivo actual: {current_file}{Colors.END}")
            print(f"{Colors.CYAN}📂 Configuracion: {config_dir}{Colors.END}")
        else:
            print(f"{Colors.YELLOW}📁 Ningun archivo seleccionado{Colors.END}")
            print(f"{Colors.CYAN}📂 Configuracion: {config_dir}{Colors.END}")
        
        opcion = input(f"\n{Colors.BOLD}Selecciona una opcion [0-8]:{Colors.END} ").strip()
        
        if opcion == '1' or opcion == '6':  # Seleccionar/cambiar archivo
            new_file = select_file()
            if new_file:
                current_file = new_file
                xonimet.set_file(current_file)
                print(f"{Colors.GREEN}✓ Archivo seleccionado: {current_file}{Colors.END}")
                input(f"\n{Colors.YELLOW}Presiona Enter para continuar...{Colors.END}")
        
        elif opcion == '2':  # Analizar archivo actual
            if not current_file:
                print(f"{Colors.RED}Primero selecciona un archivo (opcion 1){Colors.END}")
                input(f"\n{Colors.YELLOW}Presiona Enter para continuar...{Colors.END}")
                continue
            
            print(f"\n{Colors.CYAN}Analizando archivo...{Colors.END}")
            metadata = xonimet.extract_all()
            clear_screen()
            xonimet.print_metadata(metadata)
            input(f"\n{Colors.YELLOW}Presiona Enter para continuar...{Colors.END}")
        
        elif opcion == '3':  # Guardar en JSON
            if not xonimet.metadata:
                print(f"{Colors.RED}Primero analiza un archivo (opcion 2){Colors.END}")
                input(f"\n{Colors.YELLOW}Presiona Enter para continuar...{Colors.END}")
                continue
            
            save_to_json(xonimet.metadata)
            input(f"\n{Colors.YELLOW}Presiona Enter para continuar...{Colors.END}")
        
        elif opcion == '4':  # Generar PDF
            if not xonimet.metadata:
                print(f"{Colors.RED}Primero analiza un archivo (opcion 2){Colors.END}")
                input(f"\n{Colors.YELLOW}Presiona Enter para continuar...{Colors.END}")
                continue
            
            generate_pdf(xonimet)
            input(f"\n{Colors.YELLOW}Presiona Enter para continuar...{Colors.END}")
        
        elif opcion == '5':  # Ver informacion del archivo actual
            if not current_file:
                print(f"{Colors.RED}No hay archivo seleccionado{Colors.END}")
            else:
                print(f"\n{Colors.CYAN}Informacion del archivo actual:{Colors.END}")
                print(f"  {Colors.YELLOW}•{Colors.END} Ruta: {current_file}")
                print(f"  {Colors.YELLOW}•{Colors.END} Tamaño: {xonimet._format_bytes(os.path.getsize(current_file))}")
                print(f"  {Colors.YELLOW}•{Colors.END} Extension: {Path(current_file).suffix}")
                print(f"\n{Colors.CYAN}Directorio de configuracion:{Colors.END}")
                print(f"  {Colors.YELLOW}•{Colors.END} {config_dir}")
            input(f"\n{Colors.YELLOW}Presiona Enter para continuar...{Colors.END}")
        
        elif opcion == '7':  # Ayuda
            clear_screen()
            print_help()
            input(f"\n{Colors.YELLOW}Presiona Enter para continuar...{Colors.END}")
        
        elif opcion == '8':  # Limpiar pantalla
            clear_screen()
        
        elif opcion == '0':  # Salir
            print(f"\n{Colors.GREEN}Gracias por usar XONIMET 2026!{Colors.END}")
            print(f"{Colors.CYAN}Desarrollado por Darian Alberto Camacho Salas (XONIDU){Colors.END}")
            print(f"{Colors.YELLOW}Reportes guardados en: {config_dir}{Colors.END}")
            break
        
        else:
            print(f"{Colors.RED}Opcion no valida. Intenta de nuevo.{Colors.END}")
            input(f"\n{Colors.YELLOW}Presiona Enter para continuar...{Colors.END}")

# ============================================================================
# Punto de entrada
# ============================================================================
if __name__ == "__main__":
    try:
        # Modo linea de comandos
        if len(sys.argv) > 1 and sys.argv[1] not in ['-h', '--help']:
            file_path = sys.argv[1]
            if os.path.exists(file_path):
                xonimet = Xonimet(file_path)
                metadata = xonimet.extract_all()
                
                if '--json' in sys.argv:
                    print(json.dumps(metadata, indent=2, ensure_ascii=False, default=str))
                elif '--pdf' in sys.argv:
                    pdf_path = xonimet.generate_pdf_report()
                    print(f"PDF generado: {os.path.abspath(pdf_path)}")
                else:
                    xonimet.print_metadata(metadata)
            else:
                print(f"{Colors.RED}El archivo '{file_path}' no existe{Colors.END}")
        else:
            # Modo interactivo
            main()
    except KeyboardInterrupt:
        print(f"\n\n{Colors.YELLOW}Hasta pronto!{Colors.END}")
    except Exception as e:
        print(f"{Colors.RED}Error: {e}{Colors.END}")
